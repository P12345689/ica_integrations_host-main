#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Author: Mihai Criveti
Description: PowerPoint Editor Integration

This module provides a REST API endpoint for processing PowerPoint presentations
based on meeting minutes or reviewer notes, using an LLM to generate content.

The default template is based on a TDA assessment, but you can bring in your own template.

TODO:
1. Support ${} for replacements
2. Support @{20: prompt here, will result in ~20 characters, to convert to tokens} for built in prompts
3. Add support for inserting images
4. Convert chars to tokens, so say max_chars: 20 -> 5 tokens.
"""

import json
import logging
import os
import yaml
import tempfile
from typing import List, Union, Dict, Any
from datetime import date as dt_date
from uuid import uuid4

from fastapi import FastAPI, HTTPException, Request
from pydantic import BaseModel, EmailStr, Field, HttpUrl, ValidationError
from jinja2 import Environment, FileSystemLoader

from langchain_openai import AzureChatOpenAI
from langchain_community.chat_models import ChatOpenAI
from langchain_community.llms.watsonxllm import WatsonxLLM
from langchain_consultingassistants import ChatConsultingAssistants

from pptx import Presentation

from .config import get_model

# Define a Union type for different types of models
ModelType = Union[AzureChatOpenAI, ChatOpenAI, WatsonxLLM, ChatConsultingAssistants]

# Constants
DEFAULT_MAX_THREADS = int(os.getenv("DEFAULT_MAX_THREADS", "4"))
DEFAULT_TEMPLATE_URL = os.getenv("DEFAULT_TEMPLATE_URL", "app/routes/powerpoint_editor/templates/template.pptx")
DEFAULT_CONFIG_PATH = os.getenv("DEFAULT_CONFIG_PATH", "app/routes/powerpoint_editor/config.yaml")
SERVER_NAME = os.getenv("SERVER_NAME", "http://127.0.0.1:8080")

# Setup logging
debug_mode = os.getenv('DEBUG') == '1'
logging_level = logging.DEBUG if debug_mode else logging.INFO
logging.basicConfig(level=logging_level, format='%(asctime)s - %(levelname)s - %(message)s')
log = logging.getLogger(__name__)

# Load Jinja2 environment
TEMPLATES_DIRECTORY = "app/routes/powerpoint_editor/templates"
template_env = Environment(loader=FileSystemLoader(TEMPLATES_DIRECTORY))

class PowerPointEditorInput(BaseModel):
    notes: str = Field(..., description="Meeting minutes or reviewer notes")
    author_name: str = Field(None, description="Name of the author")
    author_email: EmailStr = Field(None, description="Email of the author")
    document_date: dt_date = Field(default_factory=dt_date.today, description="Date of the document")
    context: str = Field(None, description="Conversational context")
    llm_override: Union[str, List[str]] = Field(None, description="Override default LLM")
    template_url: HttpUrl = Field(None, description="URL to custom PowerPoint template")
    config_data: str = Field(None, description="YAML configuration as string")

class ReplacementInfo(BaseModel):
    placeholder: str
    value: str

class PowerPointEditorOutput(BaseModel):
    status: str = Field(default="success")
    invocationId: str
    response: List[Dict[str, str]]

def update_powerpoint(input_pptx: str, output_pptx: str, replacements: Dict[str, str]) -> None:
    """Update the PowerPoint file with the generated content."""
    prs = Presentation(input_pptx)

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        # Replace placeholders in the text
                        for placeholder, content in replacements.items():
                            if placeholder in original_text:
                                new_text = original_text.replace(placeholder, content)
                                if new_text != original_text:
                                    run.text = new_text
                                    log.info(f"Replaced '{placeholder}' with '{content[:50]}...' in shape {shape_index + 1} on slide {slide_index + 1}")

            elif shape.has_table:
                table = shape.table
                for row_index, row in enumerate(table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text
                                # Replace placeholders in the table cells
                                for placeholder, content in replacements.items():
                                    if placeholder in original_text:
                                        new_text = original_text.replace(placeholder, content)
                                        if new_text != original_text:
                                            run.text = new_text
                                            log.info(f"Replaced '{placeholder}' with '{content[:50]}...' in cell ({row_index + 1}, {cell_index + 1}) on slide {slide_index + 1}")

    # Ensure the output directory exists
    output_dir = os.path.dirname(output_pptx)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    # Save the modified presentation
    prs.save(output_pptx)
    log.info(f"Saved updated PowerPoint to {output_pptx}")

def process_powerpoint(input_data: PowerPointEditorInput, model: ModelType) -> PowerPointEditorOutput:
    """Process a PowerPoint presentation based on provided notes and configuration."""
    invocation_id = str(uuid4())
    
    # Prepare input files
    input_pptx = str(input_data.template_url or DEFAULT_TEMPLATE_URL)
    output_pptx = f"public/powerpoint/{invocation_id}.pptx"

    try:
        # Create temporary files
        with tempfile.NamedTemporaryFile(mode='w+', suffix='.yaml', delete=False) as config_file, \
             tempfile.NamedTemporaryFile(mode='w+', suffix='.txt', delete=False) as notes_file:
            
            config_path = config_file.name
            notes_path = notes_file.name

            # Write configuration to temp file if provided, else use default
            if input_data.config_data:
                config_file.write(input_data.config_data)
                config_file.flush()
            else:
                config_path = DEFAULT_CONFIG_PATH

            # Write notes and context to temp file
            notes_file.write(input_data.notes)
            if input_data.context:
                notes_file.write(f"\n\nContext:\n{input_data.context}")
            notes_file.flush()

        # Load configuration
        with open(config_path, 'r') as f:
            config = yaml.safe_load(f)

        # Load and prepare templates
        main_prompt_template = template_env.get_template("main_prompt.jinja")
        response_template = template_env.get_template("final_response.jinja")

        # Process PowerPoint
        replacements = {}
        log.info(f"Using model: {model.__class__.__name__}")
        for replacement in config['replacements']:
            prompt = main_prompt_template.render(
                main_prompt=config['main_prompt'],
                notes=input_data.notes,
                context=config.get('context', {}),
                placeholder=replacement['placeholder'],
                prompt=replacement['prompt'],
                instructions=replacement['instructions']
            )
            
            # Generate content using the LLM
            try:
                content = model.invoke([("system", "You are a helpful assistant."), ("human", prompt)])
                content = content.content if hasattr(content, 'content') else str(content)
                content = content.strip()
                log.info(f"Generated content for {replacement['placeholder']}: {content[:100]}...")
            except Exception as e:
                log.error(f"Error generating content for {replacement['placeholder']}: {str(e)}")
                content = f"Error generating content: {str(e)}"
            
            replacements[replacement['placeholder']] = content

        # Update the PowerPoint file
        try:
            update_powerpoint(input_pptx, output_pptx, replacements)
        except Exception as e:
            log.error(f"Error updating PowerPoint: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Failed to update PowerPoint: {str(e)}")
        
        pptx_url = f"{SERVER_NAME}/{output_pptx}"
        log.info(f"Processed PowerPoint: {pptx_url}")

        # Generate final response using the template
        response_content = response_template.render(
            pptx_url=pptx_url,
            replacements=[{"placeholder": k, "value": v} for k, v in replacements.items()],
            processing_date=dt_date.today().isoformat(),
            author_name=input_data.author_name,
            author_email=input_data.author_email
        )

        return PowerPointEditorOutput(
            status="success",
            invocationId=invocation_id,
            response=[
                {
                    "message": response_content,
                    "type": "text"
                }
            ]
        )

    except Exception as e:
        log.error(f"Error processing PowerPoint: {str(e)}")
        return PowerPointEditorOutput(
            status="error",
            invocationId=invocation_id,
            response=[{"message": f"Failed to process PowerPoint: {str(e)}", "type": "error"}]
        )

    finally:
        # Clean up temporary files
        if input_data.config_data and os.path.exists(config_path):
            os.remove(config_path)
        if os.path.exists(notes_path):
            os.remove(notes_path)

def add_custom_routes(app: FastAPI) -> None:
    """Adds custom routes to the FastAPI application for PowerPoint processing."""

    @app.post("/system/powerpoint_editor/process/invoke", response_model=PowerPointEditorOutput)
    async def powerpoint_editor_endpoint(request: Request) -> PowerPointEditorOutput:
        """
        Process a PowerPoint presentation based on provided notes and configuration.

        Args:
            request (Request): The request object containing the input data.

        Returns:
            PowerPointEditorOutput: The processed PowerPoint URL and list of replacements.

        Raises:
            HTTPException: If the input is invalid or processing fails.
        """
        log.info("Received request to process PowerPoint")

        try:
            data: Dict[str, Any] = await request.json()
            log.debug(f"Received data: {data}")
            input_data = PowerPointEditorInput(**data)
            log.debug(f"Validated input data: {input_data}")
        except json.JSONDecodeError:
            log.error("Invalid JSON received")
            raise HTTPException(status_code=400, detail="Invalid JSON")
        except ValidationError as e:
            log.error(f"Validation error: {e}")
            raise HTTPException(status_code=422, detail=json.loads(e.json()))

        model = get_model(input_data.llm_override)
        return process_powerpoint(input_data, model)

if __name__ == "__main__":
    # This block is for testing purposes only
    from fastapi import FastAPI
    app = FastAPI()
    add_custom_routes(app)
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8080)

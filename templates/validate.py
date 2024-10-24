#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Author: Mihai Criveti
Description: Validate pptx documents and rewrite them using python-pptx.

This script verifies if a .pptx file is corrupt, validates each element, rewrites the file in place,
and optionally saves the old version as .bak. It logs errors encountered during the process.

Requires: python-pptx

Usage:
    ./validate.py [--no-backup] <file1.pptx> <file2.pptx> ...

The script can also be imported and used as a function.

Functions:
    verify_and_rewrite_pptx(file_path: str, backup: bool = True) -> None

Arguments:
    file_path (str): The path to the .pptx file to be validated and rewritten.
    backup (bool): Whether to save the original file with a .bak extension.

Example:
    >>> verify_and_rewrite_pptx('example.pptx', backup=True)
    'example.pptx is valid. Rewritten file saved as example.pptx, original saved as example.pptx.bak'
"""

import argparse
import logging
import shutil
from typing import Optional

from pptx import Presentation

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s:%(levelname)s:%(message)s",
    handlers=[logging.FileHandler("pptx_validation.log"), logging.StreamHandler()],
)


def validate_slide(slide) -> None:
    """
    Validate each shape and text element in a slide.

    Args:
        slide: A slide object from a Presentation.

    Raises:
        ValueError: If any text in the slide is invalid.
    """
    try:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    if not isinstance(text, str):
                        raise ValueError("Invalid text in shape.")
    except Exception as e:
        logging.error(f"Error in slide validation: {e}")
        raise


def verify_and_rewrite_pptx(file_path: str, backup: bool = True) -> Optional[str]:
    """
    Verify a .pptx file is not corrupt, validate its content, and rewrite it in place.
    Optionally saves the original file with a .bak extension.

    Args:
        file_path (str): The path to the .pptx file to be validated and rewritten.
        backup (bool): Whether to save the original file with a .bak extension.

    Returns:
        Optional[str]: Message indicating the result of the operation.

    Raises:
        ValueError: If the file is not a .pptx file or if any element is invalid.
    """
    if not file_path.lower().endswith(".pptx"):
        raise ValueError("The file must be a .pptx file.")

    backup_path = file_path + ".bak"

    try:
        # Attempt to open the PowerPoint file
        prs = Presentation(file_path)
        logging.info(f"{file_path} is not corrupt.")

        # Validate each slide in the presentation
        for slide in prs.slides:
            validate_slide(slide)

        if backup:
            # Save the original file as .bak
            shutil.copy2(file_path, backup_path)

        # Save the validated presentation back to the original file
        prs.save(file_path)
        if backup:
            result_msg = f"Rewritten file saved as {file_path}, original saved as {backup_path}"
        else:
            result_msg = f"Rewritten file saved as {file_path} without creating a backup"
        logging.info(result_msg)
        return result_msg

    except Exception as e:
        # Log the error and print the error message
        logging.error(f"Error processing {file_path}: {e}")
        return f"Error: {file_path} is corrupt or contains invalid content. Exception: {e}"


def main():
    """
    Main function to process multiple .pptx files passed as command-line arguments.
    """
    parser = argparse.ArgumentParser(description="Validate and rewrite .pptx files.")
    parser.add_argument(
        "files",
        metavar="file",
        type=str,
        nargs="+",
        help="one or more .pptx files to validate and rewrite",
    )
    parser.add_argument(
        "--no-backup",
        action="store_true",
        help="do not create a .bak backup of the original file",
    )

    args = parser.parse_args()
    backup = not args.no_backup

    for file_path in args.files:
        try:
            result = verify_and_rewrite_pptx(file_path, backup)
            if result:
                print(result)
        except ValueError as ve:
            print(f"Validation error for {file_path}: {ve}")
        except Exception as e:
            print(f"An error occurred for {file_path}: {e}")


if __name__ == "__main__":
    main()
